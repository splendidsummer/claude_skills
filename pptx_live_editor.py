#!/usr/bin/env python3
"""
PPTX 实时编辑器 - 低延迟交互式编辑

使用方法:
  python pptx_live_editor.py <file.pptx>

然后输入命令:
  > edit 0 title "新标题"        # 修改第0页的标题
  > add_slide                   # 添加新幻灯片
  > list                        # 列出所有幻灯片
  > show 0                      # 显示第0页内容
  > delete 2                    # 删除第2页
  > q                           # 退出

特性:
  - 自动保存 + 自动刷新 LibreOffice
  - 增量编辑，只改需要改的部分
  - 交互式 REPL 界面
"""

import os
import sys
import subprocess
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


class PPTXLiveEditor:
    def __init__(self, filepath: str):
        self.filepath = Path(filepath).resolve()
        self.prs = Presentation(str(self.filepath))
        self.libreoffice_pid = None
        print(f"📂 已加载: {self.filepath}")
        print(f"📊 共 {len(self.prs.slides)} 页幻灯片")

    def save_and_refresh(self):
        """保存并刷新 LibreOffice"""
        self.prs.save(str(self.filepath))
        # 发送 SIGHUP 给 LibreOffice 使其刷新
        if self.libreoffice_pid:
            try:
                os.kill(self.libreoffice_pid, 1)  # SIGHUP
            except ProcessLookupError:
                pass
        print("✅ 已保存")

    def list_slides(self):
        """列出所有幻灯片"""
        print(f"\n📊 共 {len(self.prs.slides)} 页:")
        for i, slide in enumerate(self.prs.slides):
            title = self._get_slide_title(slide)
            shapes_count = len(slide.shapes)
            print(f"  [{i}] {title[:40]}... ({shapes_count} 个元素)")

    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return text
        return "(无标题)"

    def show_slide(self, index: int):
        """显示幻灯片详情"""
        if not 0 <= index < len(self.prs.slides):
            print("❌ 无效的幻灯片索引")
            return

        slide = self.prs.slides[index]
        print(f"\n📄 幻灯片 {index}:")
        print("-" * 40)

        for i, shape in enumerate(slide.shapes):
            shape_type = type(shape).__name__
            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace('\n', ' ')
                print(f"  [{i}] {shape_type}: \"{text}...\"")
            else:
                print(f"  [{i}] {shape_type}")

    def edit_text(self, slide_idx: int, search: str, new_text: str):
        """编辑文本"""
        if not 0 <= slide_idx < len(self.prs.slides):
            print("❌ 无效的幻灯片索引")
            return

        slide = self.prs.slides[slide_idx]
        found = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if search.lower() in run.text.lower():
                            # 保留原有格式
                            original_text = run.text
                            run.text = original_text.replace(search, new_text)
                            found = True
                            print(f"  ✏️ \"{original_text}\" → \"{run.text}\"")

        if found:
            self.save_and_refresh()
        else:
            print(f"❌ 未找到包含 \"{search}\" 的文本")

    def set_text(self, slide_idx: int, shape_idx: int, new_text: str):
        """直接设置指定形状的文本"""
        if not 0 <= slide_idx < len(self.prs.slides):
            print("❌ 无效的幻灯片索引")
            return

        slide = self.prs.slides[slide_idx]
        shapes = list(slide.shapes)

        if not 0 <= shape_idx < len(shapes):
            print(f"❌ 无效的形状索引 (0-{len(shapes)-1})")
            return

        shape = shapes[shape_idx]
        if shape.has_text_frame:
            shape.text_frame.paragraphs[0].runs[0].text = new_text
            print(f"  ✏️ 设置文本: \"{new_text}\"")
            self.save_and_refresh()
        else:
            print("❌ 该形状没有文本")

    def add_slide(self, layout_idx: int = 6):
        """添加新幻灯片"""
        try:
            layout = self.prs.slide_layouts[layout_idx]
            slide = self.prs.slides.add_slide(layout)
            print(f"  ➕ 添加幻灯片 [{len(self.prs.slides) - 1}]")
            self.save_and_refresh()
        except IndexError:
            print(f"❌ 无效的布局索引 (0-{len(self.prs.slide_layouts)-1})")

    def delete_slide(self, index: int):
        """删除幻灯片"""
        if not 0 <= index < len(self.prs.slides):
            print("❌ 无效的幻灯片索引")
            return

        # 获取 slide ID
        rId = self.prs.slides._sldIdLst[index].rId
        # 从关系列表中删除
        self.prs.part.drop_rel(rId)
        # 从幻灯片列表中删除
        del self.prs.slides._sldIdLst[index]

        print(f"  🗑️ 删除幻灯片 [{index}]")
        self.save_and_refresh()

    def move_slide(self, from_idx: int, to_idx: int):
        """移动幻灯片"""
        if not (0 <= from_idx < len(self.prs.slides) and 0 <= to_idx < len(self.prs.slides)):
            print("❌ 无效的幻灯片索引")
            return

        slide_id = self.prs.slides._sldIdLst[from_idx]
        del self.prs.slides._sldIdLst[from_idx]
        self.prs.slides._sldIdLst.insert(to_idx, slide_id)

        print(f"  📦 移动幻灯片 [{from_idx}] → [{to_idx}]")
        self.save_and_refresh()

    def list_layouts(self):
        """列出可用的布局"""
        print(f"\n📋 可用布局 (共 {len(self.prs.slide_layouts)} 个):")
        for i, layout in enumerate(self.prs.slide_layouts):
            name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
            print(f"  [{i}] {name}")

    def duplicate_slide(self, index: int):
        """复制幻灯片"""
        if not 0 <= index < len(self.prs.slides):
            print("❌ 无效的幻灯片索引")
            return

        source = self.prs.slides[index]
        # 使用空白布局创建新幻灯片
        blank_layout = self.prs.slide_layouts[6]  # 通常是空白布局
        new_slide = self.prs.slides.add_slide(blank_layout)

        # 复制形状
        for shape in source.shapes:
            if shape.has_text_frame:
                # 简化：只复制文本框
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                textbox.text_frame.text = shape.text_frame.text

        print(f"  📋 复制幻灯片 [{index}] → [{len(self.prs.slides) - 1}]")
        self.save_and_refresh()


def main():
    if len(sys.argv) < 2:
        print("用法: python pptx_live_editor.py <file.pptx>")
        sys.exit(1)

    editor = PPTXLiveEditor(sys.argv[1])
    print("\n💡 输入 'help' 查看命令列表\n")

    while True:
        try:
            cmd = input("\n> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n👋 再见!")
            break

        if not cmd:
            continue

        parts = cmd.split(maxsplit=2)
        action = parts[0].lower()

        try:
            if action == "q" or action == "quit" or action == "exit":
                print("👋 再见!")
                break

            elif action == "help":
                print("""
命令列表:
  list                - 列出所有幻灯片
  show <slide>        - 显示幻灯片详情
  layouts            - 列出可用布局
  edit <slide> <search> <new>  - 搜索并替换文本
  set <slide> <shape> <text>  - 直接设置形状文本
  add [layout]         - 添加新幻灯片 (默认: 6=空白)
  delete <slide>      - 删除幻灯片
  copy <slide>        - 复制幻灯片
  move <from> <to>    - 移动幻灯片
  q                  - 退出
""")

            elif action == "list" or action == "ls":
                editor.list_slides()

            elif action == "show":
                idx = int(parts[1])
                editor.show_slide(idx)

            elif action == "layouts":
                editor.list_layouts()

            elif action == "edit":
                slide_idx = int(parts[1])
                search = parts[2].strip('"')
                new_text = parts[3].strip('"')
                editor.edit_text(slide_idx, search, new_text)

            elif action == "set":
                slide_idx = int(parts[1])
                shape_idx = int(parts[2])
                text = parts[3].strip('"')
                editor.set_text(slide_idx, shape_idx, text)

            elif action == "add":
                layout_idx = int(parts[1]) if len(parts) > 1 else 6
                editor.add_slide(layout_idx)

            elif action == "delete" or action == "del":
                idx = int(parts[1])
                editor.delete_slide(idx)

            elif action == "copy" or action == "dup":
                idx = int(parts[1])
                editor.duplicate_slide(idx)

            elif action == "move" or action == "mv":
                from_idx = int(parts[1])
                to_idx = int(parts[2])
                editor.move_slide(from_idx, to_idx)

            else:
                print(f"❌ 未知命令: {action}")

        except (IndexError, ValueError) as e:
            print(f"❌ 参数错误: {e}")


if __name__ == "__main__":
    main()
