#!/usr/bin/env python3
"""
简单可靠的 PPTX 预览器 - 生成单个 HTML 文件
浏览器会自动每 2 秒重新加载页面
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("安装: pip install python-pptx")
    sys.exit(1)


def hex_color(rgb_color):
    """将 RGB 颜色转换为十六进制"""
    if rgb_color is None:
        return "#ffffff"
    return f"#{rgb_color}"


def get_text_color(shape):
    """获取形状中第一个有颜色的文本"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font and run.font.color:
                    try:
                        return hex_color(run.font.color.rgb)
                    except:
                        pass
    return "#ffffff"


def generate_html(pptx_path):
    """生成 HTML 预览文件"""
    prs = Presentation(pptx_path)

    # 幻灯片尺寸 (像素)
    width_px = 960
    height_px = 540

    slides_html = []

    for slide_idx, slide in enumerate(prs.slides):
        # 背景色
        bg_color = "#1e1e1e"
        try:
            bg = slide.background
            if bg.fill and bg.fill.type == 1:  # SOLID
                bg_color = hex_color(bg.fill.fore_color.rgb)
        except:
            pass

        shapes_html = []

        for shape in slide.shapes:
            try:
                # 位置
                left = int(shape.left * width_px / prs.slide_width) if prs.slide_width else 0
                top = int(shape.top * height_px / prs.slide_height) if prs.slide_height else 0
                width = int(shape.width * width_px / prs.slide_width) if prs.slide_width else 100
                height = int(shape.height * height_px / prs.slide_height) if prs.slide_height else 50

                # 文本
                if shape.has_text_frame:
                    text_color = get_text_color(shape)
                    text_content = []

                    for para in shape.text_frame.paragraphs:
                        para_text = ""
                        is_bold = False

                        for run in para.runs:
                            para_text += run.text
                            if run.font and run.font.bold:
                                is_bold = True

                        if para_text.strip():
                            style = f"color: {text_color};"
                            if is_bold:
                                style += " font-weight: bold;"
                            text_content.append(f'<p style="{style}">{para_text}</p>')

                    if text_content:
                        div_style = f"position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px; overflow: hidden; padding: 5px;"
                        shapes_html.append(f'<div style="{div_style}">{"".join(text_content)}</div>')

            except Exception as e:
                pass

        slide_style = f"position: relative; width: {width_px}px; height: {height_px}px; background: {bg_color}; margin: 10px auto; box-shadow: 0 4px 20px rgba(0,0,0,0.5);"
        slides_html.append(f'''
    <div class="slide" data-index="{slide_idx}" style="{slide_style}">
        {"".join(shapes_html) if shapes_html else "<p style='color: #666; text-align: center; padding-top: 200px;'>空白页</p>"}
    </div>''')

    # 完整 HTML
    html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PPTX Preview - {Path(pptx_path).name}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            background: #0a0a0a;
            color: white;
            font-family: Arial, sans-serif;
            padding: 20px;
        }}
        h1 {{
            text-align: center;
            margin-bottom: 20px;
            color: #888;
        }}
        .container {{
            max-width: 1000px;
            margin: 0 auto;
        }}
        .slide {{
            page-break-after: always;
        }}
        .nav {{
            position: fixed;
            bottom: 20px;
            right: 20px;
            background: #333;
            padding: 10px 20px;
            border-radius: 8px;
            display: flex;
            gap: 10px;
            align-items: center;
        }}
        .nav button {{
            background: #0066ff;
            border: none;
            color: white;
            padding: 8px 16px;
            border-radius: 4px;
            cursor: pointer;
        }}
        .nav button:hover {{ background: #0052cc; }}
        .nav span {{ color: #888; }}
        #current-slide {{
            display: none;
            position: fixed;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            z-index: 100;
        }}
        #current-slide.visible {{ display: block; }}
    </style>
</head>
<body>
    <h1>📊 {Path(pptx_path).name} ({len(prs.slides)} 页)</h1>
    <div class="container">
        {"".join(slides_html)}
    </div>

    <div class="nav">
        <button onclick="prevSlide()">←</button>
        <span id="page-info">1 / {len(prs.slides)}</span>
        <button onclick="nextSlide()">→</button>
        <button onclick="toggleView()">📋 全部/单页</button>
    </div>

    <div id="current-slide"></div>

    <script>
        let slides = document.querySelectorAll('.slide');
        let current = 0;
        let singleView = false;

        function showSlide(idx) {{
            if (idx < 0 || idx >= slides.length) return;
            current = idx;
            document.getElementById('page-info').textContent = (idx + 1) + ' / ' + slides.length;

            if (singleView) {{
                slides.forEach((s, i) => {{
                    s.style.display = i === idx ? 'block' : 'none';
                }});
            }}
        }}

        function nextSlide() {{ showSlide(current + 1); }}
        function prevSlide() {{ showSlide(current - 1); }}

        function toggleView() {{
            singleView = !singleView;
            if (singleView) {{
                showSlide(current);
            }} else {{
                slides.forEach(s => s.style.display = 'block');
            }}
        }}

        document.addEventListener('keydown', e => {{
            if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
            else if (e.key === 'ArrowLeft') prevSlide();
        }});

        // 自动刷新 - 每2秒检查文件修改时间
        setInterval(() => {{
            fetch(window.location.href, {{ method: 'HEAD' }})
                .then(r => {{
                    if (r.headers.get('Last-Modified')) {{
                        location.reload();
                    }}
                }});
        }}, 2000);
    </script>
</body>
</html>'''

    return html


def main():
    if len(sys.argv) < 2:
        print("用法: python pptx_simple_preview.py <file.pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"❌ 文件不存在: {pptx_path}")
        sys.exit(1)

    output_path = pptx_path.with_suffix('.preview.html')

    html = generate_html(str(pptx_path))
    output_path.write_text(html, encoding='utf-8')

    print(f"✅ 生成预览: {output_path}")
    print(f"📊 打开浏览器查看")


if __name__ == "__main__":
    main()
