#!/usr/bin/env python3
"""
PPTX Web 实时预览器 - 修复版
将 PPTX 幻灯片转换为 HTML 并在浏览器中实时预览
"""

import json
import os
import subprocess
import sys
import time
from http.server import HTTPServer, BaseHTTPRequestHandler
from pathlib import Path
from threading import Thread
import threading

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler
except ImportError:
    print("安装依赖: pip install python-pptx watchdog")
    sys.exit(1)

# 全局状态
STATE = {
    "filepath": None,
    "slides": [],
    "last_modified": 0,
    "port": 8765,
    "connected_clients": 0
}

def hex_color(rgb_color):
    """将 RGB 颜色转换为十六进制"""
    if rgb_color is None:
        return None
    # python-pptx 的 RGBColor 是 BGR 格式
    return f"#{rgb_color}"


def extract_text(shape):
    """提取形状中的文本"""
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para_text = ""
            for run in para.runs:
                text = run.text
                font = run.font
                color = None
                size = None
                bold = False

                if font:
                    try:
                        color = hex_color(font.color.rgb) if font.color and font.color.rgb else None
                    except:
                        pass
                    try:
                        size = font.size.pt if font.size else None
                    except:
                        pass
                    bold = font.bold if font.bold else False

                para_text += text

            if para_text.strip():
                alignment = str(para.alignment) if para.alignment else "left"
                texts.append({
                    "text": para_text,
                    "alignment": alignment,
                    "bold": bold
                })
    return texts


def emu_to_px(emu, total_emu, total_px):
    """将 EMU 单位转换为像素"""
    return int(emu * total_px / total_emu) if total_emu else 0


def extract_slides(filepath):
    """提取所有幻灯片内容"""
    prs = Presentation(filepath)
    slides = []

    # 幻灯片尺寸 (转换为像素，假设 96 DPI)
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    # 10英寸宽 * 96 = 960px
    slide_width_px = 960
    slide_height_px = 540

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "index": slide_idx,
            "shapes": [],
            "width": slide_width_px,
            "height": slide_height_px
        }

        # 获取背景色
        try:
            bg = slide.background
            if bg.fill:
                fill = bg.fill
                if fill.type == 1:  # SOLID
                    try:
                        slide_data["background"] = hex_color(fill.fore_color.rgb)
                    except:
                        slide_data["background"] = "#1e1e1e"
        except:
            slide_data["background"] = "#1e1e1e"

        # 提取形状
        for shape in slide.shapes:
            try:
                # 位置和大小
                left = emu_to_px(shape.left, slide_width_emu, slide_width_px)
                top = emu_to_px(shape.top, slide_height_emu, slide_height_px)
                width = emu_to_px(shape.width, slide_width_emu, slide_width_px)
                height = emu_to_px(shape.height, slide_height_emu, slide_height_px)

                shape_data = {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                    "type": str(shape.shape_type)
                }

                # 文本内容
                if shape.has_text_frame:
                    texts = extract_text(shape)
                    if texts:
                        shape_data["texts"] = texts

                        # 尝试获取文本颜色
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font and run.font.color:
                                    try:
                                        shape_data["color"] = hex_color(run.font.color.rgb)
                                    except:
                                        pass
                                break
                            break

                # 图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        shape_data["image"] = True
                        shape_data["image_bytes"] = base64.b64encode(image.blob).decode()
                        shape_data["image_ext"] = image.ext
                    except:
                        pass

                slide_data["shapes"].append(shape_data)
            except Exception as e:
                pass

        slides.append(slide_data)

    return slides


import base64


class ViewerHandler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        pass  # 静默日志

    def do_GET(self):
        if self.path == "/":
            self.send_html()
        elif self.path == "/api/slides":
            self.send_json(STATE["slides"])
        elif self.path.startswith("/api/slide/"):
            idx = int(self.path.split("/")[-1])
            if 0 <= idx < len(STATE["slides"]):
                self.send_json(STATE["slides"][idx])
            else:
                self.send_error(404)
        elif self.path == "/api/refresh":
            STATE["slides"] = extract_slides(STATE["filepath"])
            self.send_json({"status": "ok", "slides": len(STATE["slides"])})
        else:
            self.send_error(404)

    def send_html(self):
        html = '''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PPTX Preview</title>
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body {
            background: #0a0a0a;
            color: white;
            font-family: Arial, sans-serif;
            display: flex;
            flex-direction: column;
            height: 100vh;
            overflow: hidden;
        }
        #nav {
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 10px;
            background: #1a1a1a;
            gap: 20px;
        }
        #nav button {
            background: #333;
            border: none;
            color: white;
            padding: 8px 16px;
            border-radius: 4px;
            cursor: pointer;
            font-size: 14px;
        }
        #nav button:hover { background: #444; }
        #nav span { color: #888; }
        #container {
            flex: 1;
            display: flex;
            justify-content: center;
            align-items: center;
            overflow: hidden;
        }
        #slide {
            position: relative;
            background: #1e1e1e;
            box-shadow: 0 4px 20px rgba(0,0,0,0.5);
        }
        .shape {
            position: absolute;
            display: flex;
            align-items: flex-start;
            padding: 4px;
            overflow: hidden;
        }
        .shape img {
            max-width: 100%;
            max-height: 100%;
        }
    </style>
</head>
<body>
    <div id="nav">
        <button onclick="prev()">← 上一页</button>
        <span id="page">1 / 1</span>
        <button onclick="next()">下一页 →</button>
        <button onclick="refresh()">🔄 刷新</button>
        <span id="status"></span>
    </div>
    <div id="container">
        <div id="slide"></div>
    </div>

    <script>
        let slides = [];
        let current = 0;

        async function loadSlides() {
            const res = await fetch('/api/slides');
            slides = await res.json();
            show(0);
        }

        function show(idx) {
            if (idx < 0 || idx >= slides.length) return;
            current = idx;
            const slide = slides[idx];

            const el = document.getElementById('slide');
            el.style.width = slide.width + 'px';
            el.style.height = slide.height + 'px';
            el.style.background = slide.background || '#1e1e1e';
            el.innerHTML = '';

            document.getElementById('page').textContent = `${idx + 1} / ${slides.length}`;

            slide.shapes.forEach(shape => {
                if (shape.texts && shape.texts.length > 0) {
                    const div = document.createElement('div');
                    div.className = 'shape';
                    div.style.left = shape.left + 'px';
                    div.style.top = shape.top + 'px';
                    div.style.width = shape.width + 'px';
                    div.style.height = shape.height + 'px';
                    div.style.flexDirection = 'column';
                    div.style.justifyContent = 'flex-start';

                    shape.texts.forEach(t => {
                        const span = document.createElement('span');
                        span.textContent = t.text;
                        if (t.bold) span.style.fontWeight = 'bold';
                        if (shape.color) span.style.color = shape.color;
                        div.appendChild(span);
                    });

                    el.appendChild(div);
                }
            });
        }

        function next() { show(current + 1); }
        function prev() { show(current - 1); }

        async function refresh() {
            document.getElementById('status').textContent = '刷新中...';
            const res = await fetch('/api/refresh');
            const data = await res.json();
            slides = (await fetch('/api/slides')).json();
            show(current);
            document.getElementById('status').textContent = `✅ ${data.slides} 页`;
            setTimeout(() => document.getElementById('status').textContent = '', 2000);
        }

        document.addEventListener('keydown', e => {
            if (e.key === 'ArrowRight' || e.key === ' ') next();
            else if (e.key === 'ArrowLeft') prev();
        });

        loadSlides();
        setInterval(async () => {
            await refresh();
        }, 3000);
    </script>
</body>
</html>'''
        self.send_response(200)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.end_headers()
        self.wfile.write(html.encode())

    def send_json(self, data):
        self.send_response(200)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps(data).encode())

    def send_error(self, code):
        self.send_response(code)
        self.end_headers()


def main():
    if len(sys.argv) < 2:
        print("用法: python pptx_web_viewer.py <file.pptx> [port]")
        sys.exit(1)

    filepath = Path(sys.argv[1]).resolve()
    if not filepath.exists():
        print(f"❌ 文件不存在: {filepath}")
        sys.exit(1)

    STATE["filepath"] = str(filepath)
    STATE["port"] = int(sys.argv[2]) if len(sys.argv) > 2 else 8765

    print(f"📂 加载: {filepath}")
    STATE["slides"] = extract_slides(str(filepath))
    print(f"📊 共 {len(STATE['slides'])} 页")

    # 文件监听
    def on_change():
        print("🔄 检测到文件变化，重新加载...")
        STATE["slides"] = extract_slides(str(filepath))
        print(f"✅ 已更新: {len(STATE['slides'])} 页")

    class Handler(FileSystemEventHandler):
        def on_modified(self, event):
            if Path(event.src_path).resolve() == filepath:
                on_change()

    observer = Observer()
    observer.schedule(Handler(), path=str(filepath.parent), recursive=False)
    observer.start()

    # 启动服务器
    print(f"\n🌐 服务器: http://localhost:{STATE['port']}")
    print("⌨️  导航: ← → | 自动刷新: 3秒")
    print("🛑 停止: Ctrl+C\n")

    try:
        server = HTTPServer(('0.0.0.0', STATE["port"]), ViewerHandler)
        server.serve_forever()
    except KeyboardInterrupt:
        observer.stop()
        print("\n👋 已停止")

    observer.join()


if __name__ == "__main__":
    main()
